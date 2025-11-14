"""
PPT创建工具
通过JSON格式的数据生成PPT文件
"""
from langchain.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import os
from typing import Dict, List, Any

# 配色方案
COLOR_SCHEMES = {
    "blue_green": {
        "primary": RGBColor(91, 155, 213),      # 蓝色
        "secondary": RGBColor(112, 173, 171),   # 青绿色
        "accent": RGBColor(237, 125, 49),       # 橙色
        "light": RGBColor(217, 225, 242),       # 浅蓝
        "dark": RGBColor(68, 114, 196),         # 深蓝
    },
    "modern": {
        "primary": RGBColor(70, 130, 180),      # 钢蓝色
        "secondary": RGBColor(176, 196, 222),   # 浅钢蓝
        "accent": RGBColor(255, 165, 0),        # 橙色
        "light": RGBColor(240, 248, 255),       # 爱丽丝蓝
        "dark": RGBColor(25, 25, 112),          # 午夜蓝
    },
    "elegant": {
        "primary": RGBColor(95, 158, 160),      # 军服蓝
        "secondary": RGBColor(245, 222, 179),   # 小麦色
        "accent": RGBColor(210, 105, 30),       # 巧克力色
        "light": RGBColor(245, 245, 220),       # 米色
        "dark": RGBColor(47, 79, 79),           # 深石板灰
    }
}


@tool
def create_ppt_from_json(json_data: str, output_path: str = "output.pptx") -> str:
    """
    根据JSON数据创建PPT文件
    
    参数:
        json_data (str): JSON格式的PPT内容数据,包含标题、幻灯片列表等信息
            格式示例:
            {
                "title": "PPT标题",
                "slides": [
                    {
                        "type": "title",
                        "title": "主标题",
                        "subtitle": "副标题"
                    },
                    {
                        "type": "content",
                        "title": "页面标题",
                        "content": ["要点1", "要点2", "要点3"]
                    }
                ]
            }
        output_path (str): 输出的PPT文件路径,默认为"output.pptx"
    
    返回:
        str: 成功消息或错误信息
    
    异常:
        json.JSONDecodeError: JSON解析失败
        Exception: PPT创建过程中的其他异常
    """
    try:
        # 解析JSON数据
        if isinstance(json_data, str):
            data = json.loads(json_data)
        else:
            data = json_data
        
        # 确保输出目录存在
        files_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'files')
        os.makedirs(files_dir, exist_ok=True)
        
        # 如果output_path只是文件名,则添加完整路径
        if not os.path.isabs(output_path):
            output_path = os.path.join(files_dir, output_path)
        
        # 创建PPT对象
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 遍历幻灯片数据
        for slide_data in data.get("slides", []):
            slide_type = slide_data.get("type", "content")
            
            if slide_type == "title":
                # 创建标题页
                _create_title_slide(prs, slide_data)
            elif slide_type == "content":
                # 创建内容页
                _create_content_slide(prs, slide_data)
            elif slide_type == "two_column":
                # 创建双栏内容页
                _create_two_column_slide(prs, slide_data)
            elif slide_type == "catalog":
                # 创建目录页
                _create_catalog_slide(prs, slide_data)
            elif slide_type == "section":
                # 创建章节页
                _create_section_slide(prs, slide_data)
            elif slide_type == "card_grid":
                # 创建卡片网格页
                _create_card_grid_slide(prs, slide_data)
            elif slide_type == "timeline":
                # 创建时间线页
                _create_timeline_slide(prs, slide_data)
            elif slide_type == "stats":
                # 创建数据展示页
                _create_stats_slide(prs, slide_data)
            elif slide_type == "image_text":
                # 创建图文混排页
                _create_image_text_slide(prs, slide_data)
        
        # 保存PPT
        prs.save(output_path)
        return f"PPT创建成功! 文件保存在: {output_path}"
    
    except json.JSONDecodeError as e:
        return f"JSON解析错误: {str(e)}"
    except Exception as e:
        return f"PPT创建失败: {str(e)}"


def _create_title_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建标题幻灯片
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title和subtitle
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加背景装饰
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.fill.background()
    
    # 大装饰圆形
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6.5), Inches(-1.5),
        Inches(5), Inches(5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = colors["secondary"]
    circle1.line.fill.background()
    circle1.fill.transparency = 0.5
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-2), Inches(4.5),
        Inches(4.5), Inches(4.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = colors["primary"]
    circle2.line.fill.background()
    circle2.fill.transparency = 0.5
    
    # 小装饰圆点
    for pos in [(0.5, 1.5, 0.3), (8.5, 0.8, 0.25), (1.2, 6.5, 0.35)]:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(pos[0]), Inches(pos[1]),
            Inches(pos[2]), Inches(pos[2])
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors["accent"]
        dot.line.fill.background()
        dot.fill.transparency = 0.4
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "")
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = slide_data.get("subtitle", "")
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = colors["primary"]
    subtitle_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def _create_content_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建内容幻灯片
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title和content列表
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加装饰元素
    _add_decorative_shapes(slide, colors)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "")
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    
    # 内容区域
    content_list = slide_data.get("content", [])
    start_top = 2
    
    for i, item in enumerate(content_list):
        # 项目符号圆点
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(start_top + i * 0.7 + 0.1),
            Inches(0.15), Inches(0.15)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["primary"]
        bullet.line.fill.background()
        
        # 内容文字
        text_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(start_top + i * 0.7),
            Inches(8), Inches(0.6)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.color.rgb = colors["dark"]
        text_frame.word_wrap = True


def _create_two_column_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建双栏内容幻灯片
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title、left_content和right_content
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加装饰性背景元素
    _add_decorative_shapes(slide, colors)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "")
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    
    # 左栏内容框
    left_content = slide_data.get("left_content", [])
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # 右栏内容框
    right_content = slide_data.get("right_content", [])
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(12)


def _create_catalog_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建目录页
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title和items列表
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加装饰元素
    _add_decorative_shapes(slide, colors)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(7), Inches(0.8), Inches(2.5), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "目录")
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
    
    # 添加目录项
    items = slide_data.get("items", [])
    start_top = 2.2
    
    for i, item in enumerate(items):
        # 序号框
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.2), Inches(start_top + i * 0.8),
            Inches(0.5), Inches(0.5)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = colors["primary"]
        num_shape.line.color.rgb = colors["primary"]
        
        # 序号文字
        num_frame = num_shape.text_frame
        num_frame.text = f"{i+1:02d}"
        num_frame.paragraphs[0].font.size = Pt(16)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        num_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        num_frame.vertical_anchor = 1  # 居中
        
        # 目录项文字
        text_box = slide.shapes.add_textbox(
            Inches(7.9), Inches(start_top + i * 0.8),
            Inches(1.8), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.color.rgb = colors["dark"]


def _create_section_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建章节页
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含section_number和section_title
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加大背景装饰
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors["light"]
    bg_shape.line.fill.background()
    
    # 添加装饰圆形
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7), Inches(-1),
        Inches(4), Inches(4)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = colors["secondary"]
    circle1.line.fill.background()
    circle1.fill.transparency = 0.5
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1), Inches(5),
        Inches(3.5), Inches(3.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = colors["primary"]
    circle2.line.fill.background()
    circle2.fill.transparency = 0.5
    
    # 章节编号
    section_num = slide_data.get("section_number", "01")
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    num_frame = num_box.text_frame
    num_frame.text = section_num
    num_frame.paragraphs[0].font.size = Pt(120)
    num_frame.paragraphs[0].font.bold = True
    num_frame.paragraphs[0].font.color.rgb = colors["primary"]
    num_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("section_title", "")
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def _create_card_grid_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建卡片网格页
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title和cards列表
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加装饰
    _add_decorative_shapes(slide, colors)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "")
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    
    # 卡片
    cards = slide_data.get("cards", [])
    cols = 3  # 每行3个卡片
    card_width = 2.8
    card_height = 2.2
    spacing = 0.3
    start_left = 0.8
    start_top = 2
    
    for i, card in enumerate(cards[:6]):  # 最多6个卡片
        row = i // cols
        col = i % cols
        
        left = start_left + col * (card_width + spacing)
        top = start_top + row * (card_height + spacing)
        
        # 卡片背景
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(card_width), Inches(card_height)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = colors["primary"]
        card_shape.line.color.rgb = colors["dark"]
        card_shape.line.width = Pt(1)
        
        # 卡片标题
        title_text = card.get("title", "")
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.2),
            Inches(card_width - 0.4), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # 卡片内容
        content_text = card.get("content", "")
        content_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.9),
            Inches(card_width - 0.4), Inches(card_height - 1.1)
        )
        content_frame = content_box.text_frame
        content_frame.text = content_text
        content_frame.paragraphs[0].font.size = Pt(14)
        content_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        content_frame.word_wrap = True


def _create_timeline_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建时间线页
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title和events列表
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加装饰
    _add_decorative_shapes(slide, colors)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "")
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    
    # 时间线主线
    line = slide.shapes.add_connector(
        1,  # 直线
        Inches(1), Inches(3.5),
        Inches(9), Inches(3.5)
    )
    line.line.color.rgb = colors["primary"]
    line.line.width = Pt(3)
    
    # 时间点
    events = slide_data.get("events", [])
    num_events = len(events)
    
    if num_events > 0:
        spacing = 8.0 / num_events
        
        for i, event in enumerate(events):
            x_pos = 1 + (i + 0.5) * spacing
            
            # 时间点圆圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos - 0.2), Inches(3.3),
                Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = colors["accent"]
            circle.line.color.rgb = colors["dark"]
            circle.line.width = Pt(2)
            
            # 时间标签
            time_box = slide.shapes.add_textbox(
                Inches(x_pos - 0.5), Inches(2.5),
                Inches(1), Inches(0.4)
            )
            time_frame = time_box.text_frame
            time_frame.text = event.get("time", "")
            time_frame.paragraphs[0].font.size = Pt(14)
            time_frame.paragraphs[0].font.bold = True
            time_frame.paragraphs[0].font.color.rgb = colors["primary"]
            time_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # 事件描述
            desc_box = slide.shapes.add_textbox(
                Inches(x_pos - 0.6), Inches(4),
                Inches(1.2), Inches(1.5)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = event.get("description", "")
            desc_frame.paragraphs[0].font.size = Pt(12)
            desc_frame.paragraphs[0].font.color.rgb = colors["dark"]
            desc_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            desc_frame.word_wrap = True


def _create_stats_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建数据展示页
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title和stats列表
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加装饰
    _add_decorative_shapes(slide, colors)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "")
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    
    # 数据卡片
    stats = slide_data.get("stats", [])
    num_stats = len(stats)
    
    if num_stats > 0:
        card_width = 2.5
        spacing = (10 - num_stats * card_width) / (num_stats + 1)
        
        for i, stat in enumerate(stats[:4]):  # 最多4个数据
            left = spacing + i * (card_width + spacing)
            
            # 数据框
            stat_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(2.5),
                Inches(card_width), Inches(3)
            )
            stat_shape.fill.solid()
            stat_shape.fill.fore_color.rgb = colors["primary"]
            stat_shape.line.fill.background()
            
            # 数值
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(3),
                Inches(card_width - 0.4), Inches(1.2)
            )
            value_frame = value_box.text_frame
            value_frame.text = stat.get("value", "")
            value_frame.paragraphs[0].font.size = Pt(48)
            value_frame.paragraphs[0].font.bold = True
            value_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            value_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(4.3),
                Inches(card_width - 0.4), Inches(0.8)
            )
            label_frame = label_box.text_frame
            label_frame.text = stat.get("label", "")
            label_frame.paragraphs[0].font.size = Pt(16)
            label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            label_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            label_frame.word_wrap = True


def _create_image_text_slide(prs: Presentation, slide_data: Dict[str, Any]) -> None:
    """
    创建图文混排页
    
    参数:
        prs (Presentation): PPT对象
        slide_data (Dict): 幻灯片数据,包含title、content和image_placeholder
    
    返回:
        None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    colors = _get_color_scheme(slide_data.get("color_scheme", "blue_green"))
    
    # 添加装饰
    _add_decorative_shapes(slide, colors)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get("title", "")
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["dark"]
    
    # 图片占位符
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(2),
        Inches(4), Inches(4.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = colors["light"]
    img_placeholder.line.color.rgb = colors["primary"]
    img_placeholder.line.width = Pt(2)
    
    # 图片说明文字
    img_text = slide.shapes.add_textbox(
        Inches(5.7), Inches(3.5),
        Inches(3.6), Inches(1.5)
    )
    img_frame = img_text.text_frame
    img_frame.text = "[图片占位符]\n" + slide_data.get("image_description", "此处可插入图片")
    img_frame.paragraphs[0].font.size = Pt(14)
    img_frame.paragraphs[0].font.color.rgb = colors["dark"]
    img_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    
    # 文字内容
    content = slide_data.get("content", [])
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(4.5), Inches(4.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        p.font.color.rgb = colors["dark"]


def _get_color_scheme(scheme_name: str) -> Dict[str, RGBColor]:
    """
    获取配色方案
    
    参数:
        scheme_name (str): 配色方案名称
    
    返回:
        Dict[str, RGBColor]: 配色方案字典
    """
    return COLOR_SCHEMES.get(scheme_name, COLOR_SCHEMES["blue_green"])


def _add_decorative_shapes(slide, colors: Dict[str, RGBColor]) -> None:
    """
    添加装饰性图形元素
    
    参数:
        slide: 幻灯片对象
        colors (Dict): 配色方案
    
    返回:
        None
    """
    # 右上角装饰圆
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8.5), Inches(-0.5),
        Inches(2), Inches(2)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = colors["secondary"]
    circle1.line.fill.background()
    circle1.fill.transparency = 0.6
    
    # 左下角装饰圆
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.5), Inches(6.5),
        Inches(1.5), Inches(1.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = colors["primary"]
    circle2.line.fill.background()
    circle2.fill.transparency = 0.6


# 用于测试的辅助函数
def test_ppt_creation():
    """
    测试PPT创建功能,展示所有页面类型
    
    返回:
        None
    """
    test_data = {
        "title": "测试PPT",
        "slides": [
            {
                "type": "title",
                "title": "智能PPT生成系统",
                "subtitle": "基于LangGraph的专业演示文稿创建工具",
                "color_scheme": "blue_green"
            },
            {
                "type": "catalog",
                "title": "目录",
                "items": ["系统介绍", "核心功能", "技术架构", "应用案例", "未来展望"],
                "color_scheme": "blue_green"
            },
            {
                "type": "section",
                "section_number": "01",
                "section_title": "系统介绍",
                "color_scheme": "blue_green"
            },
            {
                "type": "content",
                "title": "主要功能",
                "content": [
                    "自动生成PPT结构",
                    "支持9种幻灯片类型",
                    "灵活的JSON配置",
                    "智能内容填充",
                    "美观的视觉设计"
                ],
                "color_scheme": "blue_green"
            },
            {
                "type": "card_grid",
                "title": "核心特性",
                "cards": [
                    {"title": "智能化", "content": "AI驱动的内容生成"},
                    {"title": "多样化", "content": "丰富的页面类型"},
                    {"title": "美观性", "content": "专业的视觉设计"},
                    {"title": "灵活性", "content": "高度可定制化"},
                    {"title": "高效性", "content": "快速创建PPT"},
                    {"title": "易用性", "content": "简单的操作流程"}
                ],
                "color_scheme": "blue_green"
            },
            {
                "type": "stats",
                "title": "性能数据",
                "stats": [
                    {"value": "9+", "label": "页面类型"},
                    {"value": "3", "label": "配色方案"},
                    {"value": "100%", "label": "自动化"},
                    {"value": "5分钟", "label": "创建时间"}
                ],
                "color_scheme": "blue_green"
            },
            {
                "type": "timeline",
                "title": "发展历程",
                "events": [
                    {"time": "v1.0", "description": "基础功能"},
                    {"time": "v2.0", "description": "增强设计"},
                    {"time": "v3.0", "description": "智能优化"},
                    {"time": "未来", "description": "持续创新"}
                ],
                "color_scheme": "blue_green"
            },
            {
                "type": "two_column",
                "title": "优势对比",
                "left_content": [
                    "传统方式:",
                    "• 手动创建每一页",
                    "• 耗时2-3小时",
                    "• 设计能力要求高",
                    "• 重复性劳动多"
                ],
                "right_content": [
                    "智能方式:",
                    "• AI自动生成",
                    "• 仅需5分钟",
                    "• 专业模板设计",
                    "• 一键完成创建"
                ],
                "color_scheme": "blue_green"
            },
            {
                "type": "image_text",
                "title": "应用场景",
                "content": [
                    "工作汇报: 快速生成工作总结PPT",
                    "项目展示: 专业的项目介绍演示",
                    "教学培训: 制作教学课件",
                    "商业提案: 创建商业计划书",
                    "学术报告: 生成学术演讲PPT"
                ],
                "image_description": "多场景应用示意图",
                "color_scheme": "blue_green"
            }
        ]
    }
    
    result = create_ppt_from_json.invoke({
        "json_data": json.dumps(test_data, ensure_ascii=False),
        "output_path": "test_output.pptx"
    })
    print(result)


if __name__ == "__main__":
    test_ppt_creation()
